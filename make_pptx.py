"""
Convert Financial Twin HTML slides → financial_twin.pptx
Dark glassmorphism theme matching the HTML deck.
Run: python3 make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Widescreen 16:9 ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0a, 0x0f)
BG2     = RGBColor(0x0d, 0x0d, 0x1a)
CYAN    = RGBColor(0x00, 0xd4, 0xff)
PURPLE  = RGBColor(0x7B, 0x2F, 0xBE)
GOLD    = RGBColor(0xFF, 0xD7, 0x00)
GREEN   = RGBColor(0x00, 0xe6, 0x76)
ORANGE  = RGBColor(0xff, 0x98, 0x00)
RED     = RGBColor(0xef, 0x53, 0x50)
WHITE   = RGBColor(0xe8, 0xe8, 0xf0)
MUTED   = RGBColor(0x88, 0x88, 0xa8)
GPT_C   = RGBColor(0x74, 0xaa, 0x9c)
GEM_C   = RGBColor(0x42, 0x85, 0xF4)
CLO_C   = RGBColor(0xda, 0x77, 0x56)
CARD    = RGBColor(0x12, 0x12, 0x22)

W = float(prs.slide_width)
H = float(prs.slide_height)

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide):
    """Dark background rectangle."""
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    shape.zorder = 0

def card(slide, x, y, w, h, color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color or CARD
    shape.line.color.rgb = RGBColor(0x22, 0x22, 0x38)
    shape.line.width = Pt(0.75)
    return shape

def label(slide, text, x, y, w, h, size=12, bold=False, color=None,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    return txb

def mlabel(slide, parts, x, y, w, h, size=12, align=PP_ALIGN.LEFT):
    """Multi-color text box. parts = [(text, color, bold), ...]"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    for (txt, col, bld) in parts:
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bld
        run.font.color.rgb = col
    return txb

def badge(slide, text, x, y, color=CYAN):
    b = card(slide, x, y, len(text)*0.09+0.3, 0.28, RGBColor(0x0a,0x14,0x22))
    b.line.color.rgb = color
    label(slide, text, x+0.1, y+0.02, len(text)*0.09+0.1, 0.24,
          size=7.5, bold=True, color=color)

def stage_badge(slide, text):
    b = card(slide, 0.5, 0.3, 3.5, 0.28, RGBColor(0x06,0x10,0x1a))
    b.line.color.rgb = CYAN
    label(slide, text, 0.6, 0.32, 3.3, 0.24, size=7.5, bold=True, color=CYAN)

def hline(slide, y, color=RGBColor(0x22,0x22,0x38)):
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.33), Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def bullet_block(slide, items, x, y, w, h, dot_color=CYAN, size=9.5):
    """items = [(dot_text, main_text, highlight_color_or_None)]"""
    row_h = 0.32
    for i, (dot, text, hi) in enumerate(items):
        yy = y + i * row_h
        label(slide, dot, x, yy, 0.25, row_h, size=size, bold=True, color=dot_color)
        label(slide, text, x+0.25, yy, w-0.25, row_h, size=size,
              color=hi or WHITE)

def val_row(slide, x, y, w, status_text, desc, result, status_color=GREEN):
    b = card(slide, x, y, w, 0.32, RGBColor(0x04,0x12,0x0a))
    b.line.color.rgb = RGBColor(0x00,0xe6,0x44)
    label(slide, status_text, x+0.1, y+0.05, 0.8, 0.25, size=8, bold=True, color=status_color)
    label(slide, desc, x+0.95, y+0.05, w-1.8, 0.25, size=8, color=MUTED)
    label(slide, result, x+w-0.85, y+0.05, 0.8, 0.25, size=8, color=WHITE, align=PP_ALIGN.RIGHT)

def model_card(slide, x, y, w, h, icon, name, model_str, ret, inf, top_color,
               retire="", nw30="", quote=""):
    b = card(slide, x, y, w, h)
    # top accent bar
    acc = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(3))
    acc.fill.solid(); acc.fill.fore_color.rgb = top_color; acc.line.fill.background()
    label(slide, icon, x+0.15, y+0.12, 0.5, 0.4, size=18)
    label(slide, name, x+0.15, y+0.55, w-0.3, 0.3, size=10, bold=True, color=top_color)
    label(slide, model_str, x+0.15, y+0.85, w-0.3, 0.22, size=7.5, color=MUTED)
    label(slide, f"{ret} return · {inf} inflation",
          x+0.15, y+1.07, w-0.3, 0.22, size=7.5, color=MUTED)
    if retire:
        label(slide, f"Retire: {retire}  ·  30yr: {nw30}",
              x+0.15, y+1.3, w-0.3, 0.22, size=8.5, color=WHITE, bold=True)
    if quote:
        label(slide, f'"{quote}"', x+0.15, y+1.55, w-0.3, 0.4,
              size=7.5, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)

label(s, "🧬", 0.5, 0.4, 12.33, 0.8, size=48, align=PP_ALIGN.CENTER)
label(s, "Financial Twin", 0.5, 1.2, 12.33, 1.0, size=52, bold=True,
      color=CYAN, align=PP_ALIGN.CENTER)
label(s, "Three AI Advisors. Three Visions of Your Financial Future. One Consensus.",
      0.5, 2.2, 12.33, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Model cards row
for i, (ico, name, mdl, ret, inf, col) in enumerate([
    ("⚡", "Alex Chen",    "GPT-4o",          "7.0%", "2.5%", GPT_C),
    ("🛡️","Morgan Wells", "Gemini 2.0 Flash", "5.0%", "3.5%", GEM_C),
    ("⚖️","Jordan Rivera", "Claude Sonnet 4.6","6.5%","3.0%", CLO_C),
]):
    cx = 1.3 + i * 3.7
    b = card(s, cx, 2.9, 3.3, 1.8)
    label(s, ico,  cx+0.15, 2.95, 0.5, 0.45, size=22)
    label(s, name, cx+0.15, 3.42, 3.0, 0.3, size=11, bold=True, color=col)
    label(s, mdl,  cx+0.15, 3.72, 3.0, 0.22, size=8, color=MUTED)
    label(s, f"{ret} return · {inf} inflation",
          cx+0.15, 3.94, 3.0, 0.22, size=8, color=MUTED)

# Tags
for i, t in enumerate(["MGMT690 · Spring 2026", "Jasmin Kaur",
                        "github.com/jasminkaur9/financial-twin",
                        "numpy-financial · FRED API · instructor"]):
    badge(s, t, 1.0 + i*3.0, 5.0)

label(s, "$530K divergence. Same person. Different assumptions.",
      0.5, 5.5, 12.33, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "DEFINE — 开题调研")
label(s, "The Gap No App Fills", 0.5, 0.65, 12.33, 0.7, size=34, bold=True, color=CYAN)
label(s, "Every finance app shows you the past. None simulate multiple plausible futures.",
      0.5, 1.35, 12.33, 0.35, size=11, color=MUTED)

# Left card
card(s, 0.5, 1.8, 5.9, 2.8)
label(s, "WHAT EXISTS TODAY", 0.7, 1.92, 5.5, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "Mint (RIP 2023) — showed transactions. That's it.", None),
    ("→", "YNAB / Copilot — rule-based budgeting, single projection", None),
    ("→", "Monarch Money — beautiful UI, no multi-model forecasting", None),
    ('→', 'All ask: "Where did my money go?"', GOLD),
], 0.7, 2.22, 5.5, 2.2)

# Right card
card(s, 6.6, 1.8, 6.1, 2.8)
label(s, "WHAT'S MISSING", 6.8, 1.92, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", '"Where is my money going? Across multiple plausible realities?"', GOLD),
    ("→", "Multi-LLM divergence — why do advisors disagree?", None),
    ("→", "Quantified assumptions — not one model, but a range", None),
    ("→", "FRED-grounded — real CPI, real Treasury yields", None),
], 6.8, 2.22, 5.7, 2.2)

# Research callout
card(s, 0.5, 4.75, 12.2, 0.7, RGBColor(0x03, 0x10, 0x1a))
label(s, "分头研究 (Parallel Research) confirmed: google-generativeai deprecated → switched to google-genai. "
        "Found instructor library. Found FRED API.",
      0.7, 4.88, 11.8, 0.45, size=9, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THREE-MODEL DIVERGENCE ENGINE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "REPRESENT — Roadmap")
label(s, "The Three-Model Divergence Engine", 0.5, 0.65, 12.33, 0.7, size=30, bold=True, color=CYAN)
label(s, "Not three answers — three reasoned disagreements that teach you what matters",
      0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# Assumption table
card(s, 0.5, 1.75, 12.2, 2.1)
cols = ["Advisor", "Model", "Return", "Inflation", "Philosophy"]
cw   = [2.0, 2.2, 1.2, 1.2, 5.3]
cx   = 0.6
for ci, (c, w) in enumerate(zip(cols, cw)):
    label(s, c, cx, 1.82, w, 0.25, size=7.5, bold=True, color=CYAN)
    cx += w

rows = [
    ("⚡ Alex Chen",    GPT_C, "GPT-4o",          "7.0%", "2.5%",
     "Aggressive: invest now. Market beats your debt rate."),
    ("🛡️ Morgan Wells", GEM_C, "Gemini 2.0 Flash", "5.0%", "3.5%",
     "Conservative: emergency fund first, eliminate all debt."),
    ("⚖️ Jordan Rivera", CLO_C,"Claude Sonnet 4.6","6.5%", "3.0%",
     "Evidence-based: balance investing + debt paydown."),
]
for ri, (adv, col, mdl, ret, inf, phil) in enumerate(rows):
    ry = 2.15 + ri * 0.5
    cx = 0.6
    for ci, (val, w) in enumerate(zip([adv, mdl, ret, inf, phil], cw)):
        color = col if ci == 0 else (CYAN if ci in (2,3) else WHITE)
        label(s, val, cx, ry, w, 0.4, size=9, color=color,
              bold=(ci==0))
        cx += w

# Key insight
card(s, 0.5, 3.95, 12.2, 0.9, RGBColor(0x03, 0x0d, 0x18))
label(s, "The Key Insight", 0.7, 4.02, 4, 0.25, size=8, bold=True, color=CYAN)
label(s, "Same person. Same income. Same debt. $530K difference in 30-year net worth "
        "— driven entirely by 2% return assumption spread. Making this visible is the product.",
      0.7, 4.3, 11.8, 0.45, size=9.5, color=WHITE)

# Projection lines (text representation)
card(s, 0.5, 4.95, 12.2, 0.75)
for x, txt, col in [(0.7,  "⚡ GPT-4o   ——————————————————————→  $1.40M", GPT_C),
                    (4.7,  "⚖️ Claude  ——————————————→  $1.15M",          CLO_C),
                    (8.4,  "🛡️ Gemini ——————→  $870K",                   GEM_C)]:
    label(s, txt, x, 5.1, 4.0, 0.4, size=8.5, color=col)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI COLLABORATION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "IMPLEMENT — AI as Cognition Mate")
label(s, "My Decisions. AI's Execution.", 0.5, 0.65, 12.33, 0.7, size=34, bold=True, color=CYAN)
label(s, "认知伙伴 — Neither creates alone. Both accountable.", 0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# Left — my decisions
card(s, 0.5, 1.8, 5.9, 3.5)
label(s, "MY DECISIONS — NOT THE AI'S", 0.7, 1.92, 5.5, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "Chose Personal Finance domain — simpler validation", None),
    ("→", "Decided 'Financial Twin' angle over 3 other concepts", None),
    ("→", "Crafted distinct personas + economic assumptions per model", CYAN),
    ("→", "Chose instructor after research — not AI's first suggestion", CYAN),
    ("→", "Enforced FRED ground truth as validation layer", None),
    ("→", "Modified AI outputs 12+ times based on financial logic", GOLD),
], 0.7, 2.22, 5.5, 3.0)

# Right — AI's contributions
card(s, 6.6, 1.8, 6.1, 3.5)
label(s, "WHAT THE AI CONTRIBUTED", 6.8, 1.92, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "Heavy lifting on boilerplate code (Streamlit layout)", None),
    ("→", "First drafts of Pydantic schema (I refined them)", None),
    ("→", "Research on deprecated libraries before I wrote code", None),
    ("→", "Debugging ThreadPoolExecutor timing issues", None),
    ("→", "Suggestions I rejected with reasons (audit trail shows this)", GOLD),
], 6.8, 2.22, 5.7, 3.0)

card(s, 0.5, 5.45, 12.2, 0.6, RGBColor(0x03, 0x10, 0x1a))
label(s, "Full audit trail auto-logs every prompt, response, and iteration. "
        "Downloadable as JSON from the Audit Trail tab.",
      0.7, 5.57, 11.8, 0.35, size=9, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LIVE DEMO
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "IMPLEMENT — Show Don't Tell")
label(s, "The App Running", 0.5, 0.65, 12.33, 0.7, size=34, bold=True, color=CYAN)
label(s, "jasminkaur9-financial-twin-app-npkpy9.streamlit.app",
      0.5, 1.35, 12.33, 0.3, size=11, color=CYAN, align=PP_ALIGN.CENTER)

# Profile strip
card(s, 0.5, 1.8, 12.2, 0.55, RGBColor(0x0d, 0x0d, 0x20))
for xi, (lbl, val) in enumerate([("AGE","28"),("MONTHLY INCOME","$6,500"),
                                  ("EXPENSES","$4,200"),("DEBT","$18,000 @ 5.5%")]):
    label(s, lbl, 0.7+xi*2.8, 1.85, 2.5, 0.2, size=6.5, color=MUTED, bold=True)
    label(s, val, 0.7+xi*2.8, 2.05, 2.5, 0.25, size=9.5, color=CYAN, bold=True)
card(s, 11.2, 1.85, 1.35, 0.45, RGBColor(0x06,0x18,0x2a))
label(s, "🚀 Run", 11.3, 1.95, 1.15, 0.3, size=9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Metric cards
for xi, (lbl, val) in enumerate([("CURRENT NET WORTH","$6,000"),
                                   ("SAVINGS RATE","35.4%"),
                                   ("RETIREMENT RANGE","Age 56–63"),
                                   ("30-YEAR NET WORTH","$1.15M")]):
    cx = 0.5 + xi*3.08
    b = card(s, cx, 2.45, 2.95, 0.65)
    label(s, lbl, cx+0.15, 2.5, 2.65, 0.2, size=6.5, bold=True, color=MUTED)
    label(s, val, cx+0.15, 2.72, 2.65, 0.3, size=13, bold=True, color=WHITE)

# Three model cards
for xi, (ico, name, mdl, ret, inf, col, retire, nw) in enumerate([
    ("⚡","Alex Chen","GPT-4o","7.0%","2.5%",GPT_C,"56","$1.40M"),
    ("🛡️","Morgan Wells","Gemini","5.0%","3.5%",GEM_C,"63","$870K"),
    ("⚖️","Jordan Rivera","Claude","6.5%","3.0%",CLO_C,"58","$1.15M"),
]):
    cx = 0.5 + xi * 4.08
    b = card(s, cx, 3.25, 3.9, 1.7)
    acc = s.shapes.add_shape(1, Inches(cx), Inches(3.25), Inches(3.9), Pt(3))
    acc.fill.solid(); acc.fill.fore_color.rgb = col; acc.line.fill.background()
    label(s, f"{ico} {name} · {mdl}", cx+0.15, 3.32, 3.6, 0.28, size=9, bold=True, color=col)
    label(s, f"{ret} return · {inf} inflation", cx+0.15, 3.62, 3.6, 0.22, size=7.5, color=MUTED)
    label(s, f"Retire: {retire}  ·  30yr: {nw}", cx+0.15, 3.87, 3.6, 0.25, size=9.5, bold=True, color=WHITE)

# Tags
for i, t in enumerate(["$530K divergence — same person, different assumptions",
                        "4 second parallel AI response",
                        "Full audit trail → JSON export"]):
    badge(s, t, 0.5+i*4.1, 5.15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECHNICAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "IMPLEMENT — Technical Stack")
label(s, "Expert-Level Architecture", 0.5, 0.65, 12.33, 0.7, size=34, bold=True, color=CYAN)
label(s, "Every choice made for a reason. Not just what works — what's right.",
      0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# Architecture diagram (left)
card(s, 0.5, 1.8, 5.9, 4.8, RGBColor(0x06, 0x06, 0x14))
arch_lines = [
    ("User Input", WHITE), ("  (Form / CSV Upload)", MUTED),
    ("       ↓", MUTED),
    ("numpy-financial Engine  ←── FRED API", GREEN),
    ("  FV, PV, nper             CPI · 10Y Treasury", MUTED),
    ("       ↓", MUTED),
    ("┌─────────────────────────────────┐", WHITE),
    ("│  ThreadPoolExecutor (parallel)  │", WHITE),
    ("│  GPT-4o │ Gemini 2.0 │ Claude  │  ← 4s not 18s", WHITE),
    ("│  instructor + Pydantic schema   │  ← identical output", WHITE),
    ("└─────────────────────────────────┘", WHITE),
    ("       ↓", MUTED),
    ("Synthesis + Statistical Analysis", PURPLE),
    ("  Pearson r · Spearman ρ · paired t-tests", MUTED),
    ("       ↓", MUTED),
    ("Plotly Dark Charts + Streamlit UI", CYAN),
    ("       ↓", MUTED),
    ("Audit Trail (auto-logged · JSON)", GREEN),
]
for li, (txt, col) in enumerate(arch_lines):
    label(s, txt, 0.65, 1.92+li*0.255, 5.6, 0.26, size=7.8, color=col)

# Right — decisions + CI/CD
card(s, 6.6, 1.8, 6.1, 2.4)
label(s, "KEY TECHNICAL DECISIONS", 6.8, 1.92, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "instructor — all 3 models return identical Pydantic schema", CYAN),
    ("→", "ThreadPoolExecutor — 4 seconds total, not 18 sequential", GOLD),
    ("→", "google-genai not google-generativeai — deprecated, caught early", None),
    ("→", "Demo mode — same numpy-financial engine, no API keys needed", None),
    ("→", "scipy.stats — Pearson, Spearman, paired t-tests added", GREEN),
], 6.8, 2.22, 5.7, 2.0)

card(s, 6.6, 4.35, 6.1, 2.1)
label(s, "CI/CD", 6.8, 4.47, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "GitHub Actions — lint + 34 tests on every push", None),
    ("→", "pytest — known-answer + statistical validation", CYAN),
    ("→", "Secrets scan — fails build if API key exposed", None),
], 6.8, 4.75, 5.7, 1.6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — VALIDATION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "VALIDATE — Cross-Checking Instruments")
label(s, "Trust But Verify", 0.5, 0.65, 12.33, 0.7, size=34, bold=True, color=CYAN)
label(s, "34 tests passing · FRED ground truth · Statistical significance confirmed",
      0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# Left column
label(s, "1. numpy-financial Known-Answer Tests",
      0.5, 1.8, 6.0, 0.3, size=8, bold=True, color=MUTED)
val_row(s, 0.5, 2.15, 5.9, "✅ PASS", "$18K @ 5.5% · $920/mo payment", "21.0 months")
val_row(s, 0.5, 2.55, 5.9, "✅ PASS", "$12K lump sum @ 7% · 10 years", "$23,598")
val_row(s, 0.5, 2.95, 5.9, "✅ PASS", "$500/mo annuity @ 7% · 10 years", "$86,420")
label(s, "34/34 pytest tests passing · runs on every GitHub push",
      0.5, 3.35, 5.9, 0.25, size=7.5, color=MUTED, italic=True)

label(s, "3. Statistical Significance (scipy.stats)",
      0.5, 3.7, 6.0, 0.3, size=8, bold=True, color=MUTED)
val_row(s, 0.5, 4.05, 5.9, "✅ SIG", "GPT-4o vs Gemini — Pearson r, Spearman ρ", "p < 0.001")
val_row(s, 0.5, 4.45, 5.9, "✅ SIG", "Paired t-test: GPT-4o vs Gemini trajectory", "p < 0.001")
val_row(s, 0.5, 4.85, 5.9, "✅", "Year-30 divergence — coefficient of variation", "CV ≈ 22%")
label(s, "High correlation (r ≈ 0.999) confirms divergence is structural — driven by assumptions, not noise.",
      0.5, 5.3, 5.9, 0.35, size=7.5, color=MUTED, italic=True)

# Right column
label(s, "2. FRED API Ground Truth", 6.6, 1.8, 6.1, 0.3, size=8, bold=True, color=MUTED)
card(s, 6.6, 2.15, 6.1, 2.0)
for ri, (metric, actual, rng) in enumerate([
    ("CPI Inflation",   "3.1%", "2.5% – 3.5%"),
    ("10Y Treasury",    "4.5%", "Risk premium 0.5–2.5%"),
    ("Fed Funds Rate",  "5.25%","Informs conservative outlook"),
]):
    ry = 2.25 + ri * 0.58
    label(s, metric, 6.8, ry, 2.0, 0.28, size=8.5, color=WHITE)
    label(s, actual, 8.95, ry, 1.2, 0.28, size=8.5, bold=True, color=CYAN)
    label(s, rng,   10.25, ry, 2.0, 0.28, size=7.5, color=MUTED)
    label(s, "✅",  12.25, ry, 0.3, 0.28, size=8.5, color=GREEN)

label(s, "4. Divergence Sanity Check", 6.6, 4.25, 6.1, 0.3, size=8, bold=True, color=MUTED)
card(s, 6.6, 4.6, 6.1, 1.1, RGBColor(0x04, 0x12, 0x0a))
label(s, "2% return spread × 30 years of compounding", 6.8, 4.7, 5.7, 0.28, size=8.5, color=WHITE)
label(s, "→  $530K gap on $1,300/mo investment", 6.8, 4.98, 5.7, 0.28, size=9, bold=True, color=GREEN)
label(s, "FV(7%,360,-1300,0) − FV(5%,360,-1300,0) = $524K ✅",
      6.8, 5.28, 5.7, 0.28, size=8, color=CYAN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AI LOG / AUDIT TRAIL
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "EVOLVE — Audit Trail")
label(s, "I Operated It — Here's the Evidence", 0.5, 0.65, 12.33, 0.7, size=30, bold=True, color=CYAN)
label(s, "Every prompt, every response, every iteration — auto-logged and downloadable",
      0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# Left — prompt
card(s, 0.5, 1.8, 5.9, 4.1, RGBColor(0x04, 0x04, 0x10))
label(s, "ACTUAL PROMPT SENT TO CLAUDE", 0.7, 1.92, 5.5, 0.25, size=7.5, bold=True, color=MUTED)
prompt_lines = [
    "You are Jordan Rivera, an evidence-based financial advisor.",
    "",
    "YOUR FIXED ECONOMIC ASSUMPTIONS:",
    "· Annual market return: 6.5%",
    "· Inflation rate: 3.0%",
    "",
    "Client profile:",
    "Age: 28  ·  Income: $6,500/mo",
    "Debt: $18,000 @ 5.5%  ·  Savings: $12,000",
    "",
    "Return a FinancialAnalysis Pydantic object with",
    "year-by-year projections, years 0–30...",
]
for li, line in enumerate(prompt_lines):
    col = CYAN if line.startswith("·") or line.startswith("YOUR") else \
          GOLD if any(x in line for x in ["Jordan","6.5%","3.0%"]) else \
          WHITE if line else MUTED
    label(s, line, 0.7, 2.22+li*0.245, 5.5, 0.26, size=7.8, color=col)

label(s, "Same structure sent to GPT-4o and Gemini — only assumptions differ.",
      0.7, 5.55, 5.5, 0.25, size=7.5, color=MUTED, italic=True)

# Right — iterations
card(s, 6.6, 1.8, 6.1, 4.1)
label(s, "HOW I MODIFIED THE AI — 4 KEY ITERATIONS", 6.8, 1.92, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
iters = [
    ("v1", "Models gave identical answers → Fixed by enforcing different assumption sets in each prompt", GOLD),
    ("v2", "Gemini returned unstructured text → Added Pydantic schema + response_schema parameter", CYAN),
    ("v3", "Sequential calls took 18 seconds → Switched to ThreadPoolExecutor → 4 seconds", GREEN),
    ("v4", "google-generativeai deprecated mid-build → Migrated to google-genai SDK", ORANGE),
]
for ii, (ver, txt, col) in enumerate(iters):
    iy = 2.22 + ii * 0.82
    b = card(s, 6.8, iy, 0.45, 0.35, RGBColor(0x0a,0x14,0x22))
    b.line.color.rgb = col
    label(s, ver, 6.88, iy+0.05, 0.35, 0.25, size=8, bold=True, color=col)
    label(s, txt, 7.35, iy, 5.2, 0.7, size=8.5, color=WHITE)

card(s, 6.6, 5.55, 6.1, 0.5, RGBColor(0x03, 0x10, 0x1a))
label(s, "Full audit trail auto-logs in the app → Download button exports JSON for class submission",
      6.8, 5.67, 5.7, 0.3, size=8, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — REFLECT
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s)
stage_badge(s, "REFLECT — Lessons Learned")
label(s, "The DRIVER Process in Practice", 0.5, 0.65, 12.33, 0.7, size=30, bold=True, color=CYAN)
label(s, "What the process revealed that planning alone never could",
      0.5, 1.35, 12.33, 0.3, size=11, color=MUTED)

# What worked
card(s, 0.5, 1.8, 5.9, 2.8, RGBColor(0x03, 0x10, 0x07))
label(s, "✅ WHAT WORKED", 0.7, 1.92, 5.5, 0.25, size=7.5, bold=True, color=GREEN)
bullet_block(s, [
    ("→", "分头研究 found instructor library — saved 2+ hours of JSON parsing", GREEN),
    ("→", "Research caught google-generativeai deprecation before it broke anything", None),
    ("→", "Demo mode uses same numpy-financial engine → reliable demo without API keys", None),
    ("→", "Pydantic schema enforcement made divergence analysis trivial", CYAN),
    ("→", "scipy.stats made divergence statistically verifiable (p < 0.001)", GREEN),
], 0.7, 2.22, 5.5, 2.4)

# What to change
card(s, 0.5, 4.7, 5.9, 1.55, RGBColor(0x10, 0x09, 0x02))
label(s, "⚠️ WHAT I'D CHANGE", 0.7, 4.82, 5.5, 0.25, size=7.5, bold=True, color=ORANGE)
bullet_block(s, [
    ("→", "Start with async/await instead of ThreadPoolExecutor — cleaner parallel code", None),
    ("→", "Add streaming responses — users see AI output appear in real time", None),
], 0.7, 5.1, 5.5, 1.1)

# Right column
card(s, 6.6, 1.8, 6.1, 1.5, RGBColor(0x10, 0x0e, 0x02))
label(s, "🎯 THE REAL INSIGHT", 6.8, 1.92, 5.7, 0.25, size=7.5, bold=True, color=GOLD)
label(s, "Economic assumptions are the hidden variable in every financial forecast. "
        "Making them explicit — and confirming their divergence is statistically significant — "
        "is more valuable than any single 'right' answer.",
      6.8, 2.22, 5.7, 0.9, size=9, color=WHITE, italic=True)

# Final numbers
card(s, 6.6, 3.45, 6.1, 2.8)
label(s, "FINAL NUMBERS", 6.8, 3.57, 5.7, 0.25, size=7.5, bold=True, color=MUTED)
bullet_block(s, [
    ("→", "34/34 pytest tests passing", CYAN),
    ("→", "3 frontier AI models in parallel (GPT-4o · Gemini · Claude)", None),
    ("→", "4s total AI response time (ThreadPoolExecutor)", None),
    ("→", "7 interactive Plotly charts", None),
    ("→", "$530K divergence — statistically significant (p < 0.001)", GREEN),
    ("→", "Pearson r, Spearman ρ, paired t-tests via scipy", GREEN),
], 6.8, 3.87, 5.7, 2.3)

# Links
for i, (t, col) in enumerate([
    ("github.com/jasminkaur9/financial-twin", CYAN),
    ("Live Streamlit App ↗", CYAN),
    ("Substack Post ↗", GOLD),
]):
    badge(s, t, 0.5+i*3.5, 6.55, col)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "product/financial_twin.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
